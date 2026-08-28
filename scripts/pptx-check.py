#!/usr/bin/env python3
"""PPTX 구조 검사 — decks/pptx (specs 없음, decks/pptx/README.md 참고)

이 서버엔 LibreOffice가 없어 헤드리스 렌더링 검사(check-design.py의 pptx 버전)를
할 수 없다. 대신 python-pptx로 실제로 열어 슬라이드 수·발표자 노트 존재를 확인한다 —
"파일이 손상되지 않았다"까지만 보장하고, 레이아웃은 PowerPoint에서 눈으로 봐야 한다.

사용: python3 pptx-check.py decks/w01-orientation.pptx
"""
import sys
from pptx import Presentation


def main() -> int:
    if len(sys.argv) != 2:
        print("사용: pptx-check.py <파일.pptx>")
        return 2
    path = sys.argv[1]
    try:
        p = Presentation(path)
    except Exception as e:
        print(f"✗ 열기 실패 — 파일이 손상됐을 수 있다: {e}")
        return 1

    n = len(p.slides)
    no_notes = [i for i, s in enumerate(p.slides, 1) if not s.has_notes_slide
                or not s.notes_slide.notes_text_frame.text.strip()]
    empty = [i for i, s in enumerate(p.slides, 1)
             if not any(sh.has_text_frame and sh.text_frame.text.strip() for sh in s.shapes)]

    print(f"슬라이드 {n}장")
    if no_notes:
        print(f"  발표자 노트 없음: {no_notes}")
    if empty:
        print(f"  ⚠ 텍스트 없는 슬라이드: {empty}")
    if not no_notes and not empty:
        print("  ✓ 노트·본문 전부 있음")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
